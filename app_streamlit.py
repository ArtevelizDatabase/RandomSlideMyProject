import streamlit as st
import os
import random
import io
from pptx import Presentation
import copy

# Set page config at the very beginning
st.set_page_config(page_title="SlideForge - PowerPoint Merger Pro", layout="wide")

# --- FUNGSI BANTUAN UNTUK MENYALIN SLIDE ---
def copy_slide_from_source(source_presentation, slide_index, dest_presentation, 
                          preserve_master=True, preserve_relationships=True,
                          remove_placeholder_text=True):
    """
    Menyalin sebuah slide dari satu presentasi ke presentasi lain dengan metode yang lebih aman.
    Termasuk master slides, layouts, dan graphics format yang belum diconvert.
    """
    try:
        source_slide = source_presentation.slides[slide_index]
        
        # Pastikan ukuran slide sama dengan presentasi tujuan
        if hasattr(dest_presentation, 'slide_width') and hasattr(dest_presentation, 'slide_height'):
            dest_width = dest_presentation.slide_width
            dest_height = dest_presentation.slide_height
        else:
            dest_width = source_presentation.slide_width
            dest_height = source_presentation.slide_height
            dest_presentation.slide_width = dest_width
            dest_presentation.slide_height = dest_height
        
        # Copy master slide and layout jika diaktifkan
        if preserve_master:
            dest_layout = copy_master_and_layout(source_slide, dest_presentation)
        else:
            dest_layout = get_default_layout(dest_presentation)
        
        # Buat slide baru
        dest_slide = dest_presentation.slides.add_slide(dest_layout)
        
        # Copy relationships jika diaktifkan
        if preserve_relationships:
            copy_slide_relationships(source_slide, dest_slide, source_presentation, dest_presentation)
        
        # Salin semua shapes termasuk graphics format
        copy_all_shapes(source_slide, dest_slide)
        
        # Hapus placeholder text jika diaktifkan
        if remove_placeholder_text:
            clear_placeholder_text(dest_slide)
        
        return True
        
    except Exception as e:
        print(f"Error copying slide {slide_index}: {e}")
        return False

def copy_master_and_layout(source_slide, dest_presentation):
    """Copy master slide dan layout dengan handling yang robust"""
    try:
        source_layout = source_slide.slide_layout
        source_master = source_layout.slide_master
        
        # Cari atau copy master slide
        dest_master = None
        for master in dest_presentation.slide_masters:
            if master.name == source_master.name:
                dest_master = master
                break
        
        if dest_master is None:
            # Copy master slide jika belum ada
            try:
                # Clone master slide element
                master_element = copy.deepcopy(source_master.element)
                dest_presentation.slide_masters._sldMasterLst.append(master_element)
                dest_master = dest_presentation.slide_masters[-1]
            except Exception as master_error:
                # Fallback: gunakan master default
                dest_master = dest_presentation.slide_masters[0] if dest_presentation.slide_masters else None
        
        # Cari atau copy slide layout
        dest_layout = None
        if dest_master:
            for layout in dest_master.slide_layouts:
                if layout.name == source_layout.name:
                    dest_layout = layout
                    break
        
        if dest_layout is None:
            # Copy layout jika belum ada
            try:
                if dest_master and hasattr(dest_master, 'slide_layouts'):
                    layout_element = copy.deepcopy(source_layout.element)
                    dest_master.slide_layouts._sldLayoutLst.append(layout_element)
                    dest_layout = dest_master.slide_layouts[-1]
            except Exception as layout_error:
                # Fallback ke layout yang tersedia
                if dest_master and len(dest_master.slide_layouts) > 0:
                    dest_layout = dest_master.slide_layouts[0]
                else:
                    dest_layout = get_default_layout(dest_presentation)
        
        return dest_layout or get_default_layout(dest_presentation)
        
    except Exception as e:
        return get_default_layout(dest_presentation)

def get_default_layout(dest_presentation):
    """Get default layout untuk fallback"""
    try:
        # Cari layout blank atau kosong
        for layout in dest_presentation.slide_layouts:
            if 'blank' in layout.name.lower() or 'kosong' in layout.name.lower():
                return layout
        
        # Fallback ke layout pertama
        return dest_presentation.slide_layouts[0] if dest_presentation.slide_layouts else None
    except:
        return None

def clear_placeholder_text(slide):
    """
    Menghapus teks placeholder dari slide seperti 'Click to add title', 'Click to add text', dll.
    """
    try:
        placeholder_texts = [
            "Click to add title",
            "Click to add subtitle", 
            "Click to add text",
            "Click to edit",
            "Click to add content",
            "Add your text here",
            "Type your text here",
            "Enter your title here",
            "Enter your subtitle here",
            "Title",
            "Subtitle",
            "Content placeholder",
            # Placeholder dalam bahasa Indonesia
            "Klik untuk menambahkan judul",
            "Klik untuk menambahkan subjudul",
            "Klik untuk menambahkan teks",
            "Klik untuk mengedit",
            "Ketik teks Anda di sini",
            "Masukkan judul di sini",
            "Masukkan subjudul di sini",
            "Judul",
            "Subjudul"
        ]
        
        shapes_to_check = list(slide.shapes)
        
        for shape in shapes_to_check:
            try:
                # Cek apakah shape memiliki text_frame dan text
                if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
                    current_text = shape.text.strip()
                    
                    # Hapus jika teks adalah placeholder
                    if current_text.lower() in [pt.lower() for pt in placeholder_texts]:
                        try:
                            shape.text_frame.clear()
                            # Atau set ke string kosong
                            shape.text = ""
                        except:
                            pass
                    
                    # Cek untuk placeholder yang hanya berupa teks kosong atau whitespace
                    elif not current_text or current_text.isspace():
                        try:
                            shape.text_frame.clear()
                            shape.text = ""
                        except:
                            pass
                
                # Cek untuk placeholder shapes berdasarkan format
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    try:
                        # Jika ini adalah placeholder dan tidak memiliki konten custom
                        if hasattr(shape, 'text') and shape.text.strip():
                            current_text = shape.text.strip()
                            if current_text.lower() in [pt.lower() for pt in placeholder_texts]:
                                shape.text_frame.clear()
                                shape.text = ""
                    except:
                        pass
                        
            except Exception as shape_error:
                # Skip shape yang bermasalah
                continue
                
    except Exception as e:
        print(f"Warning: Error clearing placeholder text: {e}")

def copy_slide_relationships(source_slide, dest_slide, source_presentation, dest_presentation):
    """Copy all relationships (images, media, etc.) dari source ke destination"""
    try:
        source_part = source_slide.part
        dest_part = dest_slide.part
        
        # Copy all relationships
        for rel_id, rel in source_part.rels.items():
            try:
                target_part = rel.target_part
                
                # Skip jika sudah ada relationship yang sama
                existing_rel = None
                for existing_rel_id, existing_rel in dest_part.rels.items():
                    if (hasattr(existing_rel, 'target_part') and 
                        hasattr(existing_rel.target_part, 'partname') and
                        hasattr(target_part, 'partname') and
                        existing_rel.target_part.partname == target_part.partname):
                        existing_rel = existing_rel_id
                        break
                
                if existing_rel is None:
                    # Copy part ke destination package jika belum ada
                    dest_package = dest_presentation.part.package
                    if target_part.partname not in [part.partname for part in dest_package.parts]:
                        dest_package.add_part(
                            target_part.partname,
                            target_part.content_type,
                            target_part.blob
                        )
                    
                    # Add relationship
                    dest_part.rels.add_relationship(
                        rel.reltype,
                        target_part,
                        rel_id
                    )
                        
            except Exception as rel_error:
                continue  # Skip relationship yang bermasalah
                
    except Exception as e:
        pass  # Relationship copying is optional but beneficial

def copy_all_shapes(source_slide, dest_slide):
    """Copy semua shapes termasuk graphics format yang belum diconvert"""
    try:
        # Clear existing shapes di destination slide (kecuali placeholder)
        shapes_to_remove = []
        for shape in dest_slide.shapes:
            if not (hasattr(shape, 'is_placeholder') and shape.is_placeholder):
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                dest_slide.shapes._spTree.remove(shape.element)
            except:
                pass
        
        # Copy semua shapes dari source
        for shape in source_slide.shapes:
            try:
                copy_single_shape(shape, dest_slide)
            except Exception as shape_error:
                print(f"Warning: Gagal menyalin shape: {shape_error}")
                continue
                
    except Exception as e:
        print(f"Error copying shapes: {e}")

def copy_single_shape(source_shape, dest_slide):
    """Copy single shape dengan handling khusus untuk berbagai tipe"""
    try:
        # Method 1: Deep copy element (untuk semua tipe termasuk graphics)
        if hasattr(source_shape, 'element'):
            new_element = copy.deepcopy(source_shape.element)
            dest_slide.shapes._spTree.append(new_element)
            return True
            
    except Exception as method1_error:
        try:
            # Method 2: Handle specific shape types
            shape_type = getattr(source_shape, 'shape_type', None)
            
            if shape_type == 13:  # Picture/Image
                copy_picture_shape(source_shape, dest_slide)
            elif shape_type == 19:  # Table
                copy_table_shape(source_shape, dest_slide)
            elif shape_type == 3:  # Chart
                copy_chart_shape(source_shape, dest_slide)
            elif hasattr(source_shape, 'text'):  # Text shapes
                copy_text_shape(source_shape, dest_slide)
            else:  # Generic shapes and graphics
                copy_generic_shape(source_shape, dest_slide)
                
        except Exception as method2_error:
            # Method 3: Fallback - raw element insertion
            try:
                if hasattr(source_shape, 'element'):
                    dest_slide.shapes._spTree.insert_element_before(
                        source_shape.element, 'p:extLst'
                    )
            except Exception as method3_error:
                print(f"All copy methods failed for shape: {method3_error}")
                return False
    
    return True

def copy_picture_shape(source_shape, dest_slide):
    """Copy picture/image shape"""
    try:
        if hasattr(source_shape, 'image'):
            # Get image data
            image_part = source_shape.image.part
            
            # Add image to destination
            from pptx.util import Inches
            left = getattr(source_shape, 'left', Inches(1))
            top = getattr(source_shape, 'top', Inches(1))
            width = getattr(source_shape, 'width', Inches(2))
            height = getattr(source_shape, 'height', Inches(2))
            
            dest_slide.shapes.add_picture(
                io.BytesIO(image_part.blob),
                left, top, width, height
            )
    except Exception as e:
        # Fallback to element copy
        new_element = copy.deepcopy(source_shape.element)
        dest_slide.shapes._spTree.append(new_element)

def copy_table_shape(source_shape, dest_slide):
    """Copy table shape"""
    try:
        # Tables are complex, use element copy
        new_element = copy.deepcopy(source_shape.element)
        dest_slide.shapes._spTree.append(new_element)
    except Exception as e:
        print(f"Failed to copy table: {e}")

def copy_chart_shape(source_shape, dest_slide):
    """Copy chart shape"""
    try:
        # Charts are complex, use element copy
        new_element = copy.deepcopy(source_shape.element)
        dest_slide.shapes._spTree.append(new_element)
    except Exception as e:
        print(f"Failed to copy chart: {e}")

def copy_text_shape(source_shape, dest_slide):
    """Copy text shape"""
    try:
        from pptx.util import Inches, Pt
        
        # Get position and size
        left = getattr(source_shape, 'left', Inches(1))
        top = getattr(source_shape, 'top', Inches(1))
        width = getattr(source_shape, 'width', Inches(3))
        height = getattr(source_shape, 'height', Inches(1))
        
        # Add text box
        textbox = dest_slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        
        # Copy text content and formatting
        if hasattr(source_shape, 'text_frame'):
            for paragraph in source_shape.text_frame.paragraphs:
                p = text_frame.add_paragraph()
                p.text = paragraph.text
                
                # Copy formatting if possible
                try:
                    if paragraph.font:
                        p.font.name = paragraph.font.name
                        p.font.size = paragraph.font.size
                        if hasattr(paragraph.font, 'color') and paragraph.font.color:
                            p.font.color.rgb = paragraph.font.color.rgb
                except:
                    pass
                    
    except Exception as e:
        # Fallback to element copy
        new_element = copy.deepcopy(source_shape.element)
        dest_slide.shapes._spTree.append(new_element)

def copy_generic_shape(source_shape, dest_slide):
    """Copy generic shapes including graphics format"""
    try:
        # For graphics and other complex shapes, element copy is most reliable
        new_element = copy.deepcopy(source_shape.element)
        dest_slide.shapes._spTree.append(new_element)
    except Exception as e:
        print(f"Failed to copy generic shape: {e}")

def set_standard_slide_size(presentation, width=None, height=None, size_option="16:9 Widescreen"):
    """
    Mengatur ukuran slide ke ukuran standar atau ukuran yang ditentukan.
    """
    from pptx.util import Inches
    
    if width is not None and height is not None:
        # Gunakan ukuran yang sudah ditentukan
        presentation.slide_width = width
        presentation.slide_height = height
    else:
        # Pilih ukuran berdasarkan opsi
        if size_option == "16:9 Widescreen":
            presentation.slide_width = Inches(13.333)  # 13.333 inches
            presentation.slide_height = Inches(7.5)    # 7.5 inches
        elif size_option == "4:3 Standard":
            presentation.slide_width = Inches(10)      # 10 inches
            presentation.slide_height = Inches(7.5)    # 7.5 inches
        elif size_option == "A4":
            presentation.slide_width = Inches(11.69)   # A4 width
            presentation.slide_height = Inches(8.27)   # A4 height
        else:
            # Default ke 16:9
            presentation.slide_width = Inches(13.333)
            presentation.slide_height = Inches(7.5)
    
    return presentation

# --- FUNGSI PREVIEW DAN UTILITAS ---
def extract_slide_title(slide):
    """Extract title dari slide"""
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            # Prioritaskan title shape (biasanya di posisi atas)
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                if shape.placeholder_format.type == 1:  # Title placeholder
                    return shape.text.strip()[:50]
            # Fallback ke text pertama yang ditemukan
            return shape.text.strip()[:50]
    return f"Slide Tanpa Judul"

def extract_slide_content(slide):
    """Extract all text content from slide"""
    content = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            content.append(shape.text.strip())
    return " ".join(content)[:200] + "..." if len(" ".join(content)) > 200 else " ".join(content)

def generate_slide_info(presentation):
    """Generate informasi untuk setiap slide"""
    slide_info = []
    for i, slide in enumerate(presentation.slides):
        try:
            info = {
                'index': i,
                'title': extract_slide_title(slide),
                'content': extract_slide_content(slide),
                'shapes_count': len(slide.shapes),
                'has_images': any(hasattr(shape, 'image') for shape in slide.shapes),
                'has_tables': any(shape.shape_type == 19 for shape in slide.shapes if hasattr(shape, 'shape_type')),
                'has_charts': any(shape.shape_type == 3 for shape in slide.shapes if hasattr(shape, 'shape_type'))
            }
            slide_info.append(info)
        except Exception as e:
            slide_info.append({
                'index': i,
                'title': f'Slide {i+1}',
                'content': 'Error reading slide content',
                'shapes_count': 0,
                'has_images': False,
                'has_tables': False,
                'has_charts': False
            })
    return slide_info

def show_slide_selector(uploaded_file, file_index):
    """Tampilkan UI untuk memilih slide secara manual dengan layout yang kompak"""
    try:
        prs = Presentation(io.BytesIO(uploaded_file.getvalue()))
        slide_info = generate_slide_info(prs)
        total_slides = len(slide_info)
        
        if total_slides == 0:
            st.warning(f"File tidak memiliki slide.")
            return [], "random", 1
        
        # Mode selection dalam layout yang lebih kompak
        st.write("**🎯 Pilih Mode Selection:**")
        mode_col1, mode_col2 = st.columns([2, 1])
        
        with mode_col1:
            selection_mode = st.radio(
                "Mode:",
                ["Random", "Manual", "Range"],
                key=f"mode_{file_index}_{uploaded_file.name}",
                horizontal=True,
                label_visibility="collapsed"
            )
        
        with mode_col2:
            if selection_mode == "Random":
                jumlah_ambil = st.number_input(
                    "Jumlah:",
                    min_value=1,
                    max_value=total_slides,
                    value=min(3, total_slides),
                    key=f"random_count_{file_index}_{uploaded_file.name}",
                    help=f"Pilih 1-{total_slides} slide secara random"
                )
                return [], "random", jumlah_ambil
        
        if selection_mode == "Manual":
            st.write(f"**📋 Pilih slide dari {total_slides} slide tersedia:**")
            selected_slides = []
            
            # Show slides dalam grid layout yang lebih kompak
            if total_slides <= 6:
                # Untuk file dengan sedikit slide, tampilkan semua
                cols = st.columns(min(3, total_slides))
                for idx, info in enumerate(slide_info):
                    with cols[idx % 3]:
                        with st.container():
                            st.write(f"**Slide {info['index'] + 1}**")
                            st.caption(f"{info['title'][:30]}...")
                            
                            # Indicators
                            indicators = []
                            if info['has_images']:
                                indicators.append("🖼️")
                            if info['has_tables']:
                                indicators.append("📊")
                            if info['has_charts']:
                                indicators.append("📈")
                            if indicators:
                                st.caption(" ".join(indicators))
                            
                            is_selected = st.checkbox(
                                "Pilih",
                                key=f"slide_{file_index}_{uploaded_file.name}_{info['index']}",
                                help=f"Slide {info['index'] + 1}: {info['title']}"
                            )
                            
                            if is_selected:
                                selected_slides.append(info['index'])
            else:
                # Untuk file dengan banyak slide, gunakan expander
                st.write("*Klik slide untuk melihat detail dan memilih*")
                for info in slide_info:
                    with st.expander(f"Slide {info['index'] + 1}: {info['title'][:40]}...", expanded=False):
                        detail_col1, detail_col2 = st.columns([3, 1])
                        
                        with detail_col1:
                            st.write(f"**Preview:** {info['content'][:100]}...")
                            
                            # Content indicators
                            indicators = []
                            if info['has_images']:
                                indicators.append("🖼️ Images")
                            if info['has_tables']:
                                indicators.append("📊 Tables")
                            if info['has_charts']:
                                indicators.append("📈 Charts")
                            if info['shapes_count'] > 0:
                                indicators.append(f"🔷 {info['shapes_count']} elements")
                            
                            if indicators:
                                st.caption(" • ".join(indicators))
                        
                        with detail_col2:
                            is_selected = st.checkbox(
                                "Pilih Slide",
                                key=f"slide_{file_index}_{uploaded_file.name}_{info['index']}",
                                help="Sertakan slide ini dalam presentasi final"
                            )
                            
                            if is_selected:
                                selected_slides.append(info['index'])
            
            return selected_slides, "manual", len(selected_slides)
        
        elif selection_mode == "Range":
            st.write("**📏 Pilih Rentang Slide:**")
            range_col1, range_col2 = st.columns(2)
            
            with range_col1:
                start_slide = st.number_input(
                    "Dari slide:",
                    min_value=1,
                    max_value=total_slides,
                    value=1,
                    key=f"range_start_{file_index}_{uploaded_file.name}"
                )
            
            with range_col2:
                end_slide = st.number_input(
                    "Sampai slide:",
                    min_value=start_slide,
                    max_value=total_slides,
                    value=min(start_slide + 2, total_slides),
                    key=f"range_end_{file_index}_{uploaded_file.name}"
                )
            
            range_slides = list(range(start_slide - 1, end_slide))  # Convert to 0-based index
            st.info(f"📋 Akan mengambil {len(range_slides)} slide (slide {start_slide} - {end_slide})")
            return range_slides, "range", len(range_slides)
    
    except Exception as e:
        st.error(f"Error processing file: {e}")
        return [], "random", 1

# --- FUNGSI BANTUAN UNTUK MENYALIN SLIDE ---

def add_slide_numbers(presentation, position="bottom-right"):
    """Add slide numbers ke semua slide"""
    try:
        from pptx.util import Inches, Pt
        
        for i, slide in enumerate(presentation.slides):
            try:
                # Tentukan posisi slide number
                if position == "bottom-right":
                    left = presentation.slide_width - Inches(1)
                    top = presentation.slide_height - Inches(0.5)
                elif position == "bottom-center":
                    left = presentation.slide_width / 2 - Inches(0.3)
                    top = presentation.slide_height - Inches(0.5)
                else:  # bottom-left
                    left = Inches(0.2)
                    top = presentation.slide_height - Inches(0.5)
                
                # Add text box with slide number
                textbox = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.3))
                text_frame = textbox.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = str(i + 1)
                p.font.size = Pt(12)
                p.font.name = "Arial"
                
            except Exception as shape_error:
                continue  # Skip if error adding to this slide
                
    except Exception as e:
        st.warning(f"Gagal menambahkan nomor slide: {e}")

def main():
    # Initialize session state
    if 'generation_time' not in st.session_state:
        from datetime import datetime
        st.session_state.generation_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    st.title("� SlideForge - PowerPoint Merger Pro")
    
    # Sidebar untuk settings
    with st.sidebar:
        st.header("⚙️ Advanced Settings")
        
        # Export options
        st.subheader("📤 Export Options")
        compress_output = st.checkbox("🗜️ Compress output file", value=True)
        
        # Enhanced copying options
        st.subheader("🔧 Copying Options")
        preserve_master_slides = st.checkbox("🎭 Preserve master slides", value=True, 
                                           help="Mempertahankan master slides dan layouts dari file sumber")
        preserve_relationships = st.checkbox("🔗 Preserve media relationships", value=True,
                                           help="Memastikan semua media (gambar, video) tersalin dengan benar")
        remove_placeholder_text = st.checkbox("🧹 Remove placeholder text", value=True,
                                            help="Menghapus teks placeholder dari master slides (contoh: 'Click to add title', 'Click to add text')")
    
    st.header("1. 📁 Unggah File Presentasi")
    uploaded_files = st.file_uploader(
        "Pilih file .pptx (bisa multiple files)...",
        type="pptx",
        accept_multiple_files=True,
        help="Maksimal 50MB per file. Mendukung format PowerPoint (.pptx)"
    )

    if uploaded_files:
        st.header("2. 🎯 Konfigurasi Slide Selection")
        
        # Summary info di bagian atas
        total_files = len(uploaded_files)
        total_size_mb = sum(file.size for file in uploaded_files) / (1024 * 1024)
        
        # Display summary cards
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("📁 Total Files", total_files)
        with col2:
            st.metric("💾 Total Size", f"{total_size_mb:.1f} MB")
        with col3:
            st.metric("📄 Files Valid", "Processing...")
        with col4:
            st.metric("🎯 Slides Selected", "0")
        
        st.markdown("---")
        
        # Container untuk semua file dengan scroll
        with st.container():
            # List untuk menyimpan pengaturan dari pengguna
            pengaturan_slide = []
            valid_files = []
            total_slides_selected = 0
            valid_files_count = 0
            
            # Process files dalam grid layout
            for i, uploaded_file in enumerate(uploaded_files):
                # File container dengan border - hanya expand 3 file pertama jika banyak file
                expand_default = True if total_files <= 5 else (i < 3)
                
                with st.expander(
                    f"📄 {uploaded_file.name} ({uploaded_file.size/(1024*1024):.1f} MB)", 
                    expanded=expand_default
                ):
                    try:
                        # Validasi ukuran file (maksimal 50MB)
                        if uploaded_file.size > 50 * 1024 * 1024:  # 50MB
                            st.error(f"❌ File terlalu besar (maksimal 50MB)")
                            st.error(f"Ukuran file: {uploaded_file.size/(1024*1024):.1f} MB")
                            continue
                        
                        # Membaca file dari memori untuk mendapatkan jumlah slide
                        file_bytes = uploaded_file.getvalue()
                        if len(file_bytes) == 0:
                            st.error(f"❌ File kosong atau rusak")
                            continue
                        
                        prs = Presentation(io.BytesIO(file_bytes))
                        total_slides = len(prs.slides)
                        
                        if total_slides == 0:
                            st.warning(f"⚠️ File tidak memiliki slide")
                            continue
                        
                        # File info dalam layout yang kompak
                        info_col1, info_col2, info_col3 = st.columns([2, 1, 1])
                        with info_col1:
                            st.write(f"**Status:** ✅ Valid PowerPoint")
                        with info_col2:
                            st.metric("Slides", total_slides, delta=None)
                        with info_col3:
                            st.metric("Size", f"{uploaded_file.size/(1024*1024):.1f} MB", delta=None)
                        
                        # Slide selector dalam container yang rapi
                        with st.container():
                            selected_slides, mode, count = show_slide_selector(uploaded_file, i)
                            
                            if mode == "random" and count > 0:
                                # Generate random slides sekarang untuk konsistensi
                                try:
                                    temp_prs = Presentation(io.BytesIO(uploaded_file.getvalue()))
                                    total_slide_sumber = len(temp_prs.slides)
                                    if total_slide_sumber > 0:
                                        list_index_slide = list(range(total_slide_sumber))
                                        random_slides = random.sample(list_index_slide, min(count, total_slide_sumber))
                                        pengaturan_slide.append((uploaded_file, random_slides, "random", count))
                                        slide_numbers = [str(idx + 1) for idx in sorted(random_slides)]
                                        st.success(f"✅ {count} slide random dipilih: {', '.join(slide_numbers)}")
                                    else:
                                        pengaturan_slide.append((uploaded_file, [], "random", count))
                                        st.success(f"✅ {count} slide random akan diambil")
                                except:
                                    pengaturan_slide.append((uploaded_file, [], "random", count))
                                    st.success(f"✅ {count} slide random akan diambil")
                                total_slides_selected += count
                            elif mode in ["manual", "range"] and selected_slides:
                                pengaturan_slide.append((uploaded_file, selected_slides, mode, len(selected_slides)))
                                total_slides_selected += len(selected_slides)
                                st.success(f"✅ {len(selected_slides)} slide dipilih")
                            elif mode in ["manual", "range"]:
                                st.warning(f"⚠️ Tidak ada slide yang dipilih")
                            
                        valid_files.append(uploaded_file)
                        valid_files_count += 1
                        
                    except Exception as e:
                        st.error(f"❌ Gagal membaca file: {str(e)}")
                        st.error("Pastikan file adalah format PowerPoint (.pptx) yang valid")
            
            # Update summary metrics
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("📁 Total Files", total_files, delta=None)
            with col2:
                st.metric("💾 Total Size", f"{total_size_mb:.1f} MB", delta=None)
            with col3:
                st.metric("📄 Files Valid", valid_files_count, delta=valid_files_count)
            with col4:
                st.metric("🎯 Slides Selected", total_slides_selected, delta=total_slides_selected)

        # Hanya tampilkan section proses jika ada file valid
        if pengaturan_slide:
            st.header("3. ⚙️ Konfigurasi Output")
            
            # Show summary of selections
            total_selected = sum(item[3] for item in pengaturan_slide)
            st.success(f"✅ Total {total_selected} slide akan digabungkan dari {len(pengaturan_slide)} file")
            
            # Detail breakdown slide yang akan diambil
            with st.expander("📋 Detail Slide yang Akan Diambil", expanded=True):
                slide_counter = 1
                for i, (source_file, selected_slides, mode, count) in enumerate(pengaturan_slide):
                    st.write(f"**📄 File: {source_file.name}**")
                    
                    if mode == "random":
                        st.info(f"🎲 Mode Random: {count} slide akan dipilih secara acak dari file ini")
                        
                        # Generate random slide selection untuk preview
                        try:
                            # Baca file untuk mendapatkan total slide
                            file_bytes = source_file.getvalue()
                            temp_prs = Presentation(io.BytesIO(file_bytes))
                            total_slide_sumber = len(temp_prs.slides)
                            
                            if total_slide_sumber > 0:
                                # Pilih slide random untuk preview
                                list_index_slide = list(range(total_slide_sumber))
                                random_slides = random.sample(list_index_slide, min(count, total_slide_sumber))
                                random_slides_sorted = sorted(random_slides)
                                
                                # Tampilkan slide yang dipilih
                                slide_numbers = [str(idx + 1) for idx in random_slides_sorted]
                                st.write(f"**Slide yang dipilih:** {', '.join(slide_numbers)}")
                                
                                # Detail mapping slide
                                st.write("**Mapping slide:**")
                                for idx in random_slides_sorted:
                                    st.write(f"• Slide #{slide_counter} ← Slide {idx + 1} dari {source_file.name}")
                                    slide_counter += 1
                            else:
                                st.warning("File tidak memiliki slide")
                                
                        except Exception as e:
                            st.warning("Tidak dapat membaca file untuk preview")
                            # Fallback ke preview umum
                            for j in range(count):
                                st.write(f"• Slide #{slide_counter + j} (Random dari {source_file.name})")
                            slide_counter += count
                        
                    elif mode == "manual" and selected_slides:
                        slide_numbers = [str(idx + 1) for idx in sorted(selected_slides)]
                        st.info(f"✋ Mode Manual: Slide nomor {', '.join(slide_numbers)} dipilih")
                        
                        # Detail mapping slide
                        st.write("**Mapping slide:**")
                        for idx in sorted(selected_slides):
                            st.write(f"• Slide #{slide_counter} ← Slide {idx + 1} dari {source_file.name}")
                            slide_counter += 1
                            
                    elif mode == "range" and selected_slides:
                        first_slide = min(selected_slides) + 1
                        last_slide = max(selected_slides) + 1
                        st.info(f"📏 Mode Range: Slide {first_slide} sampai {last_slide}")
                        
                        # Detail mapping slide
                        st.write("**Mapping slide:**")
                        for idx in sorted(selected_slides):
                            st.write(f"• Slide #{slide_counter} ← Slide {idx + 1} dari {source_file.name}")
                            slide_counter += 1
                    
                    if i < len(pengaturan_slide) - 1:  # Tidak tampilkan divider di file terakhir
                        st.markdown("---")
                
                st.success(f"🎯 **Total Slide Hasil Akhir: {slide_counter - 1} slide**")
            
            # Pilihan ukuran slide dan nama file
            col1, col2 = st.columns(2)
            with col1:
                ukuran_slide = st.selectbox(
                    "📐 Ukuran slide untuk hasil akhir:",
                    ["Auto (dari file pertama)", "16:9 Widescreen", "4:3 Standard", "A4"],
                    index=0,
                    key="slide_size_selection"
                )
            
            with col2:
                nama_file_output = st.text_input(
                    "📁 Nama file hasil:", 
                    "Presentasi_Gabungan_Pro.pptx",
                    help="File akan otomatis berakhiran .pptx"
                )

            # Tampilkan info ukuran yang dipilih
            size_info = {
                "16:9 Widescreen": "📏 13.33\" × 7.5\" (standar presentasi modern)",
                "4:3 Standard": "📏 10\" × 7.5\" (standar presentasi klasik)", 
                "A4": "📏 11.69\" × 8.27\" (format dokumen A4)",
                "Auto (dari file pertama)": "📏 Ukuran mengikuti file pertama"
            }
            if ukuran_slide in size_info:
                st.info(size_info[ukuran_slide])

            if st.button("🚀 Gabungkan Presentasi!", type="primary", use_container_width=True):
                # Validasi input
                if not pengaturan_slide:
                    st.error("Tidak ada file yang valid untuk diproses.")
                    return
                
                if not nama_file_output.strip():
                    st.error("Nama file output tidak boleh kosong.")
                    return
                
                # Pastikan nama file berakhiran .pptx
                if not nama_file_output.lower().endswith('.pptx'):
                    nama_file_output += '.pptx'
                
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                with st.spinner('🔄 Memproses presentasi Anda...'):
                    try:
                        presentasi_final = Presentation()
                        total_slides_processed = 0
                        target_width = None
                        target_height = None
                        
                        # Update progress
                        progress_bar.progress(10)
                        status_text.text("📐 Menentukan ukuran slide...")
                        
                        # Tentukan ukuran slide berdasarkan pilihan user
                        if ukuran_slide == "Auto (dari file pertama)":
                            # Ambil ukuran dari file pertama yang valid
                            for source_file, _, _, _ in pengaturan_slide:
                                try:
                                    temp_stream = io.BytesIO(source_file.getvalue())
                                    temp_prs = Presentation(temp_stream)
                                    if len(temp_prs.slides) > 0:
                                        target_width = temp_prs.slide_width
                                        target_height = temp_prs.slide_height
                                        break
                                except:
                                    continue
                            
                            # Set ukuran dari file pertama atau default jika gagal
                            if target_width and target_height:
                                presentasi_final = set_standard_slide_size(presentasi_final, target_width, target_height)
                                st.info(f"📐 Menggunakan ukuran dari file pertama: {target_width/914400:.1f}\" x {target_height/914400:.1f}\"")
                            else:
                                presentasi_final = set_standard_slide_size(presentasi_final, size_option="16:9 Widescreen")
                                st.info("📐 Menggunakan ukuran default: 16:9 Widescreen")
                        else:
                            # Gunakan ukuran yang dipilih user
                            presentasi_final = set_standard_slide_size(presentasi_final, size_option=ukuran_slide)
                            target_width = presentasi_final.slide_width
                            target_height = presentasi_final.slide_height
                            st.info(f"📐 Menggunakan ukuran {ukuran_slide}: {target_width/914400:.1f}\" x {target_height/914400:.1f}\"")
                        
                        # Process each file
                        total_files = len(pengaturan_slide)
                        
                        for file_index, (source_file, selected_slides, mode, count) in enumerate(pengaturan_slide):
                            progress = 20 + (file_index / total_files) * 60
                            progress_bar.progress(int(progress))
                            status_text.text(f"📄 Memproses '{source_file.name}'... ({file_index + 1}/{total_files})")
                            
                            try:
                                # Membaca file yang diunggah dari memori
                                source_stream = io.BytesIO(source_file.getvalue())
                                presentasi_sumber = Presentation(source_stream)
                                total_slide_sumber = len(presentasi_sumber.slides)

                                if total_slide_sumber == 0:
                                    st.warning(f"'{source_file.name}' tidak memiliki slide. Melewati file ini.")
                                    continue
                                
                                # Cek ukuran slide sumber
                                source_width = presentasi_sumber.slide_width
                                source_height = presentasi_sumber.slide_height
                                
                                if source_width != target_width or source_height != target_height:
                                    st.warning(f"⚠️ '{source_file.name}' memiliki ukuran slide berbeda. Slide akan disesuaikan ke ukuran target.")
                                
                                # Determine which slides to copy
                                if mode == "random":
                                    if selected_slides:  # Jika sudah ada slide yang dipilih dari preview
                                        index_slide_terpilih = selected_slides
                                    else:  # Fallback jika belum ada
                                        list_index_slide = list(range(total_slide_sumber))
                                        index_slide_terpilih = random.sample(list_index_slide, min(count, total_slide_sumber))
                                else:  # manual or range
                                    index_slide_terpilih = selected_slides
                                
                                st.write(f"-> Mengambil {len(index_slide_terpilih)} slide dari {total_slide_sumber} slide.")

                                successful_copies = 0
                                for index in sorted(index_slide_terpilih): # diurutkan agar lebih rapi
                                    if copy_slide_from_source(
                                        presentasi_sumber, index, presentasi_final,
                                        preserve_master=preserve_master_slides,
                                        preserve_relationships=preserve_relationships,
                                        remove_placeholder_text=remove_placeholder_text
                                    ):
                                        successful_copies += 1
                                    else:
                                        st.warning(f"Gagal menyalin slide {index + 1} dari '{source_file.name}'")
                                
                                total_slides_processed += successful_copies
                                st.success(f"✅ Berhasil menyalin {successful_copies} slide dari '{source_file.name}'")
                                
                            except Exception as file_error:
                                st.error(f"❌ Error memproses file '{source_file.name}': {file_error}")
                                continue

                        if total_slides_processed == 0:
                            st.error("❌ Tidak ada slide yang berhasil diproses. Silakan periksa file Anda.")
                            return

                        # Apply enhancements
                        progress_bar.progress(85)
                        status_text.text("🎨 Menerapkan enhancements...")
                        
                        # No template or slide numbering applied - keeping it simple

                        # Save final presentation
                        progress_bar.progress(95)
                        status_text.text("💾 Menyimpan presentasi final...")
                        
                        final_presentation_stream = io.BytesIO()
                        presentasi_final.save(final_presentation_stream)
                        final_presentation_stream.seek(0)

                        progress_bar.progress(100)
                        status_text.text("✅ Selesai!")
                        
                        st.success(f"🎉 Presentasi berhasil digabungkan dengan {total_slides_processed} slide!")
                        
                        # Enhanced download section
                        st.download_button(
                            label="📥 Unduh Presentasi (.pptx)",
                            data=final_presentation_stream,
                            file_name=nama_file_output,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            type="primary",
                            use_container_width=True
                        )
                        
                        # Show final statistics
                        st.markdown("### 📊 Statistik Final")
                        col1, col2, col3 = st.columns(3)
                        
                        with col1:
                            st.metric("Total Slide", total_slides_processed)
                        with col2:
                            st.metric("File Diproses", len(pengaturan_slide))
                        with col3:
                            st.metric("Ukuran", ukuran_slide)
                        
                    except Exception as e:
                        st.error(f"❌ Terjadi kesalahan saat memproses: {e}")
                        st.error("Silakan coba lagi atau pastikan file PowerPoint Anda tidak rusak.")
                        
        else:
            st.warning("📋 Belum ada file yang valid dipilih. Silakan upload file PowerPoint (.pptx) terlebih dahulu.")
    
    else:
        st.info("� Silakan unggah file presentasi untuk memulai proses penggabungan.")

# Menjalankan fungsi utama
if __name__ == "__main__":
    main()
